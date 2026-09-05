from __future__ import annotations

import html
import json
import math
from pathlib import Path

import pandas as pd
from docx import Document
from flask import current_app
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches as PptInches, Pt

from ..core.database import Database, utcnow


def _db() -> Database:
    return current_app.extensions["meridian_db"]


def _frame_from_payload(payload: dict, workspace_id: str) -> tuple[pd.DataFrame, str]:
    result_id = payload.get("result_id")
    if result_id:
        result = _db().get("query_results", str(result_id))
        if not result or result.get("workspace_id", "default") != workspace_id:
            raise ValueError("查询结果不存在或不属于当前工作空间")
        return pd.read_csv(result["path"]), str(result.get("sql", ""))
    rows = payload.get("rows")
    if rows is None and (payload.get("sections") or payload.get("slides")):
        return pd.DataFrame(), str(payload.get("sql", ""))
    if not isinstance(rows, list):
        raise ValueError("需要 result_id 或 rows")
    return pd.DataFrame(rows), str(payload.get("sql", ""))


def _safe_spreadsheet_text(value):
    if not isinstance(value, str):
        return value
    stripped = value.lstrip(" \t\r\n")
    if stripped.startswith(("=", "+", "-", "@")):
        return "'" + value
    return value


def _safe_spreadsheet_frame(frame: pd.DataFrame) -> pd.DataFrame:
    safe = frame.copy()
    safe.columns = [_safe_spreadsheet_text(str(column)) for column in safe.columns]
    for column in safe.columns:
        if pd.api.types.is_object_dtype(safe[column]) or pd.api.types.is_string_dtype(safe[column]):
            safe[column] = safe[column].map(_safe_spreadsheet_text)
    return safe


COLOR_SCHEMES = {
    "mckinsey": ["#003B71", "#005CAB", "#0083CA", "#00A3E0", "#7FBA00", "#FFC000"],
    "bcg": ["#006C5B", "#009879", "#00B398", "#CDECE5", "#A6192E", "#999999"],
    "bain": ["#E41E26", "#FF5C5C", "#A6192E", "#F4E8E9", "#00B398", "#999999"],
    "ey": ["#FFD100", "#FFED70", "#75787B", "#D9D9D6", "#7FBA00", "#DA3B01"],
}


def _rgb(value: str) -> RGBColor:
    clean = value.strip().lstrip("#")
    if len(clean) == 3:
        clean = "".join(character * 2 for character in clean)
    if len(clean) != 6 or any(character not in "0123456789abcdefABCDEF" for character in clean):
        raise ValueError(f"无效颜色：{value}")
    return RGBColor(int(clean[:2], 16), int(clean[2:4], 16), int(clean[4:], 16))


def _render_ppt_outline(payload: dict, path: Path) -> int:
    from ..reference_output.PPT import constants
    from ..reference_output.PPT import MckEngine

    slides = payload.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("PPT 大纲需要至少一张幻灯片")
    if len(slides) > 100:
        raise ValueError("PPT 最多 100 张幻灯片")
    supported = {
        "cover", "toc", "section_divider", "closing", "big_number", "two_stat",
        "metric_cards", "data_table", "table_insight", "executive_summary",
        "two_column_text", "action_items", "donut", "grouped_bar", "stacked_bar", "timeline",
    }
    scheme = payload.get("color_scheme") if isinstance(payload.get("color_scheme"), dict) else {}
    scheme_name = str(scheme.get("name") or payload.get("color_scheme") or "mckinsey").lower()
    colors = scheme.get("colors") or COLOR_SCHEMES.get(scheme_name) or COLOR_SCHEMES["mckinsey"]
    fallback = [_rgb(value) for value in colors]
    symbolic = {
        name: getattr(constants, name)
        for name in (
            "NAVY", "ACCENT_BLUE", "ACCENT_GREEN", "ACCENT_ORANGE", "ACCENT_RED",
            "WHITE", "BLACK", "DARK_GRAY", "MED_GRAY", "LINE_GRAY", "BG_GRAY",
            "LIGHT_BLUE", "LIGHT_GREEN", "LIGHT_ORANGE", "LIGHT_RED",
        )
    }
    symbolic.update({
        "PRIMARY": fallback[0], "SECONDARY": fallback[1 % len(fallback)],
        "ACCENT": fallback[2 % len(fallback)], "POSITIVE": fallback[3 % len(fallback)],
        "NEGATIVE": fallback[4 % len(fallback)], "NEUTRAL": fallback[5 % len(fallback)],
    })

    def resolve(value):
        if isinstance(value, str):
            if value.strip().upper() in symbolic:
                return symbolic[value.strip().upper()]
            if value.strip().startswith("#"):
                return _rgb(value)
            return value
        if isinstance(value, list):
            return [resolve(item) for item in value]
        if isinstance(value, dict):
            return {key: resolve(item) for key, item in value.items()}
        return value

    template_path = Path(__file__).resolve().parents[1] / "reference_output" / "PPT" / "PPT_template" / f"{scheme_name}.pptx"
    engine = MckEngine(total_slides=len(slides), template=str(template_path) if template_path.is_file() else None)
    for index, slide in enumerate(slides, 1):
        if not isinstance(slide, dict):
            raise ValueError(f"第 {index} 张幻灯片必须是对象")
        layout = str(slide.get("layout") or "")
        if layout not in supported:
            raise ValueError(f"第 {index} 张使用不支持的布局：{layout}")
        params = resolve(slide.get("params") or {})
        if layout in {"grouped_bar", "stacked_bar"} and "series" in params:
            normalized = []
            for series_index, item in enumerate(params["series"]):
                color = fallback[series_index % len(fallback)]
                if isinstance(item, dict):
                    normalized.append((str(item.get("name") or f"Series {series_index + 1}"), resolve(item.get("color")) if item.get("color") else color))
                elif isinstance(item, (list, tuple)):
                    normalized.append((str(item[0]), resolve(item[1]) if len(item) > 1 else color))
                else:
                    normalized.append((str(item), color))
            params["series"] = normalized
        if layout == "metric_cards" and "cards" in params:
            normalized_cards = []
            for card_index, item in enumerate(params["cards"]):
                if isinstance(item, dict):
                    normalized_cards.append((
                        str(item.get("letter") or card_index + 1), str(item.get("card_title") or item.get("title") or ""),
                        str(item.get("description") or item.get("desc") or ""), fallback[card_index % len(fallback)],
                        constants.BG_GRAY,
                    ))
                else:
                    normalized_cards.append(tuple(item) if isinstance(item, list) else (str(card_index + 1), str(item), ""))
            params["cards"] = normalized_cards
        if layout == "table_insight":
            params.setdefault("insights", ["—"])
        if layout == "executive_summary":
            params.setdefault("headline", "")
            params.setdefault("items", [])
        if layout == "action_items":
            params.setdefault("actions", [])
        if layout == "two_column_text":
            params.setdefault("columns", [])
        try:
            getattr(engine, layout)(**params)
        except Exception as exc:
            raise ValueError(f"第 {index} 张幻灯片（{layout}）生成失败：{exc}") from exc
    engine.save(str(path))
    return len(slides)


def _artifact(path: Path, kind: str, workspace_id: str, title: str, metadata: dict | None = None) -> dict:
    return _db().put(
        "artifacts",
        {
            "id": _db().new_id("art"),
            "workspace_id": workspace_id,
            "title": title,
            "kind": kind,
            "filename": path.name,
            "path": str(path),
            "size": path.stat().st_size,
            "metadata": metadata or {},
            "status": "ready",
            "created_at": utcnow(),
        },
        workspace_id=workspace_id,
    )


def export_data(payload: dict, workspace_id: str) -> dict:
    frames = payload.get("frames") if isinstance(payload.get("frames"), dict) else None
    if frames:
        normalized_frames = {
            str(name)[:80]: value if isinstance(value, pd.DataFrame) else pd.DataFrame(value)
            for name, value in frames.items()
        }
        frame = next(iter(normalized_frames.values()))
        sql = str(payload.get("sql") or "")
    else:
        frame, sql = _frame_from_payload(payload, workspace_id)
        normalized_frames = {"分析结果": frame}
    kind = str(payload.get("format", "xlsx")).lower()
    title = str(payload.get("title") or "分析结果")[:100]
    artifact_id = _db().new_id("export")
    if kind == "csv":
        path = current_app.config["SETTINGS"].export_dir / f"{artifact_id}.csv"
        _safe_spreadsheet_frame(frame).to_csv(path, index=False, encoding="utf-8-sig")
    elif kind == "xlsx":
        path = current_app.config["SETTINGS"].export_dir / f"{artifact_id}.xlsx"
        with pd.ExcelWriter(path, engine="openpyxl") as writer:
            used_names = set()
            for raw_name, current in normalized_frames.items():
                base = "".join("_" if character in "[]:*?/\\" else character for character in raw_name)[:31] or "Sheet"
                sheet_name = base
                suffix = 2
                while sheet_name in used_names:
                    tail = f"_{suffix}"
                    sheet_name = base[:31 - len(tail)] + tail
                    suffix += 1
                used_names.add(sheet_name)
                _safe_spreadsheet_frame(current).to_excel(writer, sheet_name=sheet_name, index=False)
                sheet = writer.book[sheet_name]
                header_fill = PatternFill("solid", fgColor="14213D")
                for cell in sheet[1]:
                    cell.font = Font(color="FFFFFF", bold=True)
                    cell.fill = header_fill
                    cell.alignment = Alignment(vertical="center")
                sheet.freeze_panes = "A2"
                sheet.auto_filter.ref = sheet.dimensions
                for column in sheet.columns:
                    width = min(48, max(12, max(len(str(cell.value or "")) for cell in column) + 2))
                    sheet.column_dimensions[column[0].column_letter].width = width
            if sql:
                meta = writer.book.create_sheet("分析口径")
                meta.append(["生成时间", utcnow()])
                meta.append(["只读查询", sql])
                meta.column_dimensions["A"].width = 18
                meta.column_dimensions["B"].width = 100
    else:
        raise ValueError("数据导出格式必须是 csv 或 xlsx")
    return _artifact(path, kind, workspace_id, title, {
        "rows": sum(len(value) for value in normalized_frames.values()),
        "tables": list(normalized_frames), "columns": list(frame.columns), "sql": sql,
    })


def export_report(payload: dict, workspace_id: str) -> dict:
    frame, sql = _frame_from_payload(payload, workspace_id)
    kind = str(payload.get("format", "docx")).lower()
    title = str(payload.get("title") or "数据分析报告")[:100]
    summary = str(payload.get("summary") or "本报告由经纬分析工作台根据已执行的只读查询生成。")
    insights = payload.get("insights") if isinstance(payload.get("insights"), list) else []
    sections = payload.get("sections") if isinstance(payload.get("sections"), list) else []
    artifact_id = _db().new_id("report")
    if kind == "docx":
        path = current_app.config["SETTINGS"].export_dir / f"{artifact_id}.docx"
        document = Document()
        document.add_heading(title, 0)
        document.add_paragraph(f"生成时间：{utcnow()}")
        document.add_heading("执行摘要", level=1)
        document.add_paragraph(summary)
        for section in sections:
            if not isinstance(section, dict):
                continue
            heading = str(section.get("heading") or "")
            content = section.get("content", "")
            if heading:
                document.add_heading(heading, level=1)
            if isinstance(content, list) and content and isinstance(content[0], dict):
                columns = list(content[0])
                section_table = document.add_table(rows=1, cols=max(1, len(columns)))
                section_table.style = "Table Grid"
                for index, column in enumerate(columns):
                    section_table.rows[0].cells[index].text = str(column)
                for row in content:
                    cells = section_table.add_row().cells
                    for index, column in enumerate(columns):
                        cells[index].text = str(row.get(column, ""))[:500]
            else:
                for paragraph in str(content).split("\n\n"):
                    if paragraph.strip():
                        document.add_paragraph(paragraph.strip())
        if insights:
            document.add_heading("关键洞察", level=1)
            for item in insights:
                document.add_paragraph(str(item), style="List Bullet")
        document.add_heading("数据明细", level=1)
        shown = frame.head(100)
        table = document.add_table(rows=1, cols=max(1, len(shown.columns)))
        table.style = "Light Shading Accent 1"
        for index, column in enumerate(shown.columns):
            table.rows[0].cells[index].text = str(column)
        for row in shown.itertuples(index=False):
            cells = table.add_row().cells
            for index, value in enumerate(row):
                cells[index].text = "" if pd.isna(value) else str(value)[:200]
        if sql:
            document.add_heading("可复核查询", level=1)
            document.add_paragraph(sql)
        document.save(path)
    elif kind == "pptx":
        path = current_app.config["SETTINGS"].export_dir / f"{artifact_id}.pptx"
        if payload.get("slides"):
            slide_count = _render_ppt_outline(payload, path)
            return _artifact(path, kind, workspace_id, title, {
                "rows": len(frame), "sql": sql, "slides": slide_count,
                "color_scheme": payload.get("color_scheme") or "mckinsey",
            })
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"经纬分析工作台 · {utcnow()[:10]}"
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = "执行摘要"
        slide.placeholders[1].text = summary
        if insights:
            slide = deck.slides.add_slide(deck.slide_layouts[1])
            slide.shapes.title.text = "关键洞察"
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            for index, item in enumerate(insights[:8]):
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                paragraph.text = str(item)
                paragraph.level = 0
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        slide.shapes.title.text = "数据样本"
        shown = frame.head(12)
        rows, cols = len(shown) + 1, min(len(shown.columns), 8)
        table = slide.shapes.add_table(rows, max(cols, 1), PptInches(0.5), PptInches(1.4), PptInches(12.2), PptInches(5)).table
        for col in range(cols):
            table.cell(0, col).text = str(shown.columns[col])
        for row_index, row in enumerate(shown.itertuples(index=False), 1):
            for col in range(cols):
                value = row[col]
                table.cell(row_index, col).text = "" if pd.isna(value) else str(value)[:80]
        for cell in table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
        deck.save(path)
    else:
        raise ValueError("报告格式必须是 docx 或 pptx")
    return _artifact(path, kind, workspace_id, title, {
        "rows": len(frame), "sql": sql, "sections": len(sections),
    })


def export_dashboard_html(dashboard: dict, workspace_id: str) -> dict:
    artifact_id = _db().new_id("board")
    path = current_app.config["SETTINGS"].export_dir / f"{artifact_id}.html"
    widgets = dashboard.get("widgets", []) if isinstance(dashboard.get("widgets"), list) else []
    cards = [
        (
            f'<article class="widget"><h2>{html.escape(str(widget.get("title") or "图表"))}</h2>'
            f'<div class="chart" id="widget-{index}" role="img" '
            f'aria-label="{html.escape(str(widget.get("title") or "数据图表"), quote=True)}"></div></article>'
        )
        for index, widget in enumerate(widgets)
    ]
    vendor_dir = Path(__file__).resolve().parents[2] / "frontend" / "vendor"
    vendor_path = vendor_dir / "echarts.min.js"
    map_path = vendor_dir / "echarts-china.min.js"
    if not vendor_path.is_file() or not map_path.is_file():
        raise FileNotFoundError("ECharts 离线资源不存在")
    echarts_source = vendor_path.read_text(encoding="utf-8").replace("</script", "<\\/script")
    map_source = map_path.read_text(encoding="utf-8").replace("</script", "<\\/script")

    def json_safe(value):
        if isinstance(value, dict):
            return {str(key): json_safe(item) for key, item in value.items()}
        if isinstance(value, (list, tuple)):
            return [json_safe(item) for item in value]
        if isinstance(value, float) and not math.isfinite(value):
            return None
        return value

    widget_json = json.dumps(json_safe(widgets), ensure_ascii=False, default=str).replace("<", "\\u003c")
    content = f"""<!doctype html><html lang="zh-CN"><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1"><title>{html.escape(dashboard.get('name', '分析看板'))}</title>
<style>body{{font:16px system-ui;margin:0;background:#eef2f7;color:#14213d}}header{{padding:32px 5vw;background:#14213d;color:#fff}}main{{display:grid;grid-template-columns:repeat(auto-fit,minmax(340px,1fr));gap:20px;padding:24px 5vw}}article{{background:#fff;border-radius:14px;padding:20px;box-shadow:0 8px 24px #14213d12}}h2{{font-size:16px;margin:0 0 8px}}.chart{{height:360px}}.kpi{{display:grid;place-content:center;text-align:center;font-size:42px;font-weight:700}}.kpi small{{font-size:14px;color:#667085}}.error{{display:grid;place-content:center;color:#b42318}}@media print{{body{{background:white}}article{{break-inside:avoid;box-shadow:none;border:1px solid #ddd}}}}</style>
<header><h1>{html.escape(dashboard.get('name', '分析看板'))}</h1><p>{html.escape(dashboard.get('description', ''))}</p></header><main>{''.join(cards)}</main>
<script>{echarts_source}</script><script>{map_source}</script><script id="dashboard-data" type="application/json">{widget_json}</script>
<script>(function(){{
const palette=['#167c80','#e59b4c','#4058b4','#9a5bc4','#42a46f','#df6b62'];
const widgets=JSON.parse(document.getElementById('dashboard-data').textContent);
const fallback=spec=>({{tooltip:{{trigger:'axis'}},legend:{{bottom:0}},grid:{{left:48,right:24,top:35,bottom:55,containLabel:true}},xAxis:{{type:'category',data:spec.encoding?.x||[]}},yAxis:{{type:'value'}},series:(spec.encoding?.series||[]).map(item=>({{name:item.name,type:['line','area','stacked_area'].includes(spec.type)?'line':'bar',data:item.values,areaStyle:['area','stacked_area'].includes(spec.type)?{{opacity:.2}}:undefined,stack:spec.type==='stacked_area'||spec.type==='stacked_bar'?'total':undefined}}))}});
widgets.forEach((widget,index)=>{{
  const root=document.getElementById(`widget-${{index}}`); if(!root)return;
  if(widget.error||widget.refresh_status==='error'){{root.classList.add('error');root.textContent=widget.error||widget.refresh_error||'组件刷新失败';return;}}
  if(widget.type==='kpi'||widget.chart_type==='KPI_Card'){{root.classList.add('kpi');root.innerHTML='';const value=document.createElement('div');value.textContent=widget.kpi_value??'—';const sub=document.createElement('small');sub.textContent=[widget.kpi_sub,widget.kpi_trend==null?'':`${{widget.kpi_trend>0?'↑':'↓'}} ${{Math.abs(widget.kpi_trend)}}%`].filter(Boolean).join(' · ');root.append(value,sub);return;}}
  const spec=widget.chart||widget.spec||{{}};const option=spec.option||fallback(spec);option.color=option.color||palette;option.animation=false;echarts.init(root).setOption(option,true);
}});
addEventListener('resize',()=>document.querySelectorAll('.chart').forEach(root=>echarts.getInstanceByDom(root)?.resize()));
}})();</script></html>"""
    path.write_text(content, encoding="utf-8")
    return _artifact(path, "html", workspace_id, dashboard.get("name", "分析看板"), {"dashboard_id": dashboard.get("id")})
